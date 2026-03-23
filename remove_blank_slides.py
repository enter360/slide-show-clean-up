#!/usr/bin/env python3
"""
remove_blank_slides.py
======================
Remove blank slides from .pptx and blank pages from .pdf files.

Usage:
    python remove_blank_slides.py input.pptx [--output output.pptx] [--threshold 0.01]
    python remove_blank_slides.py input.pdf  [--output output.pdf]  [--threshold 0.01]

Dependencies:
    pip install python-pptx pypdf pdfplumber Pillow
"""

import argparse
import sys
from pathlib import Path


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------

def slide_is_blank(slide, text_threshold: int = 3) -> bool:
    """
    Return True if a PPTX slide has no meaningful content.

    A slide is considered blank when ALL of the following are true:
      - No shape has visible text longer than `text_threshold` characters
      - No shape contains an image / picture
      - No shape has a non-default fill (solid colour, gradient, pattern)
    """
    from pptx.util import Pt
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.dml.color import RGBColor
    import pptx.oxml.ns as ns

    for shape in slide.shapes:
        # ── images ──────────────────────────────────────────────────────────
        if shape.shape_type == 13:          # MSO_SHAPE_TYPE.PICTURE
            return False

        # pictures embedded as placeholders
        if hasattr(shape, "image"):
            return False

        # ── text ────────────────────────────────────────────────────────────
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if len(text) > text_threshold:
                return False

        # ── non-trivial fills on any shape ──────────────────────────────────
        try:
            fill = shape.fill
            # fill.type is None for "no fill" / inherited
            if fill.type is not None:
                # Solid white / transparent fills are still "blank"
                if fill.type.name == "SOLID":
                    try:
                        rgb = fill.fore_color.rgb
                        # Treat pure-white fills as blank
                        if rgb != RGBColor(0xFF, 0xFF, 0xFF):
                            return False
                    except Exception:
                        return False          # can't determine → assume content
                else:
                    return False              # gradient, pattern, etc.
        except Exception:
            pass

        # ── grouped shapes (recursive) ───────────────────────────────────────
        if shape.shape_type == 6:           # MSO_SHAPE_TYPE.GROUP
            # If a group exists with children, treat as content
            if len(shape.shapes) > 0:
                return False

    return True


def remove_blank_slides_pptx(input_path: Path, output_path: Path, text_threshold: int = 3) -> dict:
    """Remove blank slides from a PPTX file and write the result."""
    from pptx import Presentation
    from pptx.oxml.ns import qn
    import copy

    prs = Presentation(str(input_path))
    xml_slides = prs.slides._sldIdLst          # the XML element holding slide refs

    total = len(prs.slides)
    blank_indices = []

    for i, slide in enumerate(prs.slides):
        if slide_is_blank(slide, text_threshold):
            blank_indices.append(i)

    # Remove in reverse order so indices stay valid
    for i in sorted(blank_indices, reverse=True):
        slide = prs.slides[i]
        rId = prs.slides._sldIdLst[i].get("r:id")
        # Remove the relationship and the slide part
        prs.part.drop_rel(rId)
        del xml_slides[i]

    prs.save(str(output_path))

    return {
        "total": total,
        "removed": len(blank_indices),
        "kept": total - len(blank_indices),
        "blank_slide_numbers": [i + 1 for i in blank_indices],
    }


# ---------------------------------------------------------------------------
# PDF helpers
# ---------------------------------------------------------------------------

def pdf_page_is_blank(page, text_threshold: int = 3, image_check: bool = True) -> bool:
    """
    Return True if a pdfplumber page has no meaningful content.

    Checks:
      - Extracted text length ≤ text_threshold
      - No embedded images (when image_check=True)
      - No drawn lines / curves / rectangles
    """
    # ── text ────────────────────────────────────────────────────────────────
    text = page.extract_text() or ""
    if len(text.strip()) > text_threshold:
        return False

    # ── images ──────────────────────────────────────────────────────────────
    if image_check and page.images:
        return False

    # ── vector graphics (lines, curves, rects) ──────────────────────────────
    if page.lines or page.curves or page.rects:
        return False

    return True


def remove_blank_pages_pdf(input_path: Path, output_path: Path, text_threshold: int = 3) -> dict:
    """Remove blank pages from a PDF file and write the result."""
    import pdfplumber
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(str(input_path))
    total = len(reader.pages)
    blank_indices = []

    with pdfplumber.open(str(input_path)) as pdf:
        for i, page in enumerate(pdf.pages):
            if pdf_page_is_blank(page, text_threshold):
                blank_indices.append(i)

    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        if i not in blank_indices:
            writer.add_page(page)

    # Preserve metadata
    if reader.metadata:
        writer.add_metadata(reader.metadata)

    with open(str(output_path), "wb") as f:
        writer.write(f)

    return {
        "total": total,
        "removed": len(blank_indices),
        "kept": total - len(blank_indices),
        "blank_page_numbers": [i + 1 for i in blank_indices],
    }


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="Remove blank slides (.pptx) or blank pages (.pdf) from a file.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    p.add_argument("input", help="Path to the input .pptx or .pdf file")
    p.add_argument(
        "-o", "--output",
        default=None,
        help=(
            "Path for the output file. "
            "Defaults to <input_stem>_cleaned.<ext> in the same directory."
        ),
    )
    p.add_argument(
        "--text-threshold",
        type=int,
        default=3,
        metavar="N",
        help=(
            "Maximum number of non-whitespace characters a slide/page may contain "
            "before it is considered non-blank. Default: 3."
        ),
    )
    p.add_argument(
        "--dry-run",
        action="store_true",
        help="Report which slides/pages would be removed without writing output.",
    )
    return p


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()

    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        sys.exit(f"Error: file not found – {input_path}")

    suffix = input_path.suffix.lower()
    if suffix not in {".pptx", ".pdf"}:
        sys.exit(f"Error: unsupported file type '{suffix}'. Must be .pptx or .pdf.")

    if args.output:
        output_path = Path(args.output).expanduser().resolve()
    else:
        output_path = input_path.with_name(f"{input_path.stem}_cleaned{suffix}")

    print(f"Input : {input_path}")
    print(f"Output: {output_path}")
    print(f"Text threshold: {args.text_threshold} chars\n")

    if args.dry_run:
        print("DRY RUN – no file will be written.\n")

    # ── process ─────────────────────────────────────────────────────────────
    try:
        if suffix == ".pptx":
            if not args.dry_run:
                result = remove_blank_slides_pptx(input_path, output_path, args.text_threshold)
            else:
                # For dry-run: detect only
                from pptx import Presentation
                prs = Presentation(str(input_path))
                total = len(prs.slides)
                blank_indices = [
                    i for i, s in enumerate(prs.slides)
                    if slide_is_blank(s, args.text_threshold)
                ]
                result = {
                    "total": total,
                    "removed": len(blank_indices),
                    "kept": total - len(blank_indices),
                    "blank_slide_numbers": [i + 1 for i in blank_indices],
                }

            label_total  = "slides"
            label_blank  = "blank_slide_numbers"

        else:  # .pdf
            if not args.dry_run:
                result = remove_blank_pages_pdf(input_path, output_path, args.text_threshold)
            else:
                import pdfplumber
                with pdfplumber.open(str(input_path)) as pdf:
                    total = len(pdf.pages)
                    blank_indices = [
                        i for i, p in enumerate(pdf.pages)
                        if pdf_page_is_blank(p, args.text_threshold)
                    ]
                result = {
                    "total": total,
                    "removed": len(blank_indices),
                    "kept": total - len(blank_indices),
                    "blank_page_numbers": [i + 1 for i in blank_indices],
                }

            label_total = "pages"
            label_blank = "blank_page_numbers"

    except ImportError as e:
        sys.exit(
            f"Missing dependency: {e}\n"
            "Install with:  pip install python-pptx pypdf pdfplumber Pillow"
        )

    # ── report ───────────────────────────────────────────────────────────────
    print(f"Total {label_total}  : {result['total']}")
    print(f"Blank {label_total}  : {result['removed']}")
    print(f"Kept  {label_total}  : {result['kept']}")

    blank_nums = result.get(label_blank, [])
    if blank_nums:
        print(f"Blank {label_total[:-1]} numbers: {blank_nums}")
    else:
        print(f"No blank {label_total} found.")

    if not args.dry_run and result["removed"] > 0:
        print(f"\n✓ Saved cleaned file to: {output_path}")
    elif not args.dry_run:
        print(f"\n✓ No changes needed – file is already clean.")


if __name__ == "__main__":
    main()